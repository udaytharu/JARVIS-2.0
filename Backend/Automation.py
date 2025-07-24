import requests  # Added to fix NameError
from AppOpener import close, open as appopen
from pywhatkit import search as pywhatkit_search, playonyt
from dotenv import dotenv_values
from bs4 import BeautifulSoup
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import webbrowser
import subprocess
import keyboard
import asyncio
import os
import logging
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import pyautogui


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler()]
)

# Configure pyautogui safety settings
pyautogui.FAILSAFE = True  # Move mouse to corner to abort
pyautogui.PAUSE = 0.1  # Small pause between actions

# Load environment variables with detailed error handling
env_vars = {}
try:
    with open(".env", "r", encoding="utf-8") as f:
        for i, line in enumerate(f, 1):
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" not in line:
                logging.warning(f"Invalid .env line {i}: {line}")
                continue
            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip()
            if not key or not value:
                logging.warning(f"Empty key or value in .env line {i}: {line}")
                continue
            env_vars[key] = value
except Exception as e:
    logging.error(f"Failed to read .env file: {str(e)}")

GROQ_API_KEY = env_vars.get("GroqAPIKey")
EMAIL_ADDRESS = env_vars.get("EmailAddress")
EMAIL_PASSWORD = env_vars.get("EmailPassword")

# Initialize Groq client
client = None
if GROQ_API_KEY:
    try:
        client = Groq(api_key=GROQ_API_KEY)
        client.models.list()  # Test API connectivity
        logging.info("Successfully initialized Groq client")
    except Exception as e:
        logging.error(f"Groq initialization failed: {str(e)}")
else:
    logging.warning("GROQ_API_KEY not found or invalid in .env file. AI features disabled.")

# Configure HTTP session
try:
    session = requests.Session()
    session.headers.update({
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    })
    logging.info("HTTP session initialized successfully")
except Exception as e:
    logging.error(f"Failed to initialize HTTP session: {str(e)}")
    session = None

# AI Configuration
SYSTEM_PROMPT = {
    "role": "system",
    "content": "You are a professional writing assistant. Generate high-quality content in clear, concise English."
}
messages = [SYSTEM_PROMPT]

def GoogleSearch(query: str) -> bool:
    """Perform Google search"""
    try:
        pywhatkit_search(query)
        logging.info(f"Google search executed for: {query}")
        return True
    except Exception as e:
        logging.error(f"Google search failed: {str(e)}")
        return False

def Content(topic: str) -> bool:
    """Generate and save AI content using Groq"""
    if not client:
        logging.error("Cannot generate content: Groq client not initialized")
        return False

    def generate_content(prompt: str) -> str:
        """Generate content using Groq"""
        try:
            messages.append({"role": "user", "content": prompt})
            response = client.chat.completions.create(
                model="llama3-70b-8192",
                messages=messages[-6:],
                temperature=0.7,
                max_tokens=2000,
                top_p=1.0
            )
            content = response.choices[0].message.content.strip()
            messages.append({"role": "assistant", "content": content})
            logging.info(f"Generated content for prompt: {prompt[:50]}...")
            return content
        except Exception as e:
            logging.error(f"Content generation failed: {str(e)}")
            return f"Error: {str(e)}"

    def save_and_open(content: str, filename: str) -> bool:
        """Save content to file and open in editor"""
        try:
            os.makedirs("data", exist_ok=True)
            filepath = os.path.join("data", f"{filename}.txt")
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            logging.info(f"Saved content to {filepath}")
            editor = "notepad.exe" if os.name == "nt" else "gedit"
            subprocess.Popen([editor, filepath])
            return True
        except Exception as e:
            logging.error(f"File operation failed: {str(e)}")
            return False

    logging.info(f"Starting content generation for: {topic}")
    clean_topic = topic.strip()[:100]
    content = generate_content(clean_topic)
    
    if content.startswith("Error:"):
        logging.error(content)
        return False
        
    filename = clean_topic.lower().replace(" ", "_").replace(":", "").replace("/", "_")
    return save_and_open(f"Topic: {clean_topic}\n\n{content}", filename)

def CreateGammaPresentation(topic: str) -> bool:
    """Generate a PowerPoint presentation using Groq"""
    if not client:
        logging.error("Cannot generate presentation: Groq client not initialized")
        return False

    def generate_presentation_content(prompt: str) -> str:
        """Generate presentation content using Groq"""
        try:
            messages.append({"role": "user", "content": f"Create a presentation outline about {prompt} with 6 slides: Title, Purpose, Key Features, How It Works, Example Usage, and Conclusion. For each slide, provide a title and 3-5 concise bullet points. Format as plain text with slide titles prefixed by 'Slide X: ' and bullet points prefixed by '- '. Separate slides with a blank line."})
            response = client.chat.completions.create(
                model="llama3-70b-8192",
                messages=messages[-6:],
                temperature=0.7,
                max_tokens=2000,
                top_p=1.0
            )
            content = response.choices[0].message.content.strip()
            messages.append({"role": "assistant", "content": content})
            logging.info(f"Generated presentation content for: {prompt[:50]}...")
            return content
        except Exception as e:
            logging.error(f"Presentation content generation failed: {str(e)}")
            return f"Error: {str(e)}"

    def save_and_open_presentation(content: str, filename: str) -> bool:
        """Save presentation content to PowerPoint file and open"""
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout to avoid placeholder issues

            # Parse content into slides
            slides = content.split("\n\n")
            for slide_text in slides:
                lines = slide_text.split("\n")
                if not lines or not lines[0].startswith("Slide"):
                    continue

                # Create a new slide
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title text box
                title_text = lines[0].replace("Slide ", "").replace(": ", "")
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(32)

                # Add content text box
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                text_frame = content_box.text_frame
                for bullet in lines[1:]:
                    if bullet.startswith("- "):
                        p = text_frame.add_paragraph()
                        p.text = bullet[2:].strip()
                        p.level = 0
                        p.font.size = Pt(18)
                        p.alignment = PP_ALIGN.LEFT

            # Save the presentation
            os.makedirs("data", exist_ok=True)
            filepath = os.path.join("data", f"{filename}.pptx")
            prs.save(filepath)
            logging.info(f"Presentation saved to {filepath}")

            # Open the presentation
            if os.name == "nt":
                os.startfile(filepath)  # Windows
            else:
                subprocess.Popen(["xdg-open", filepath])  # Linux
            return True
        except Exception as e:
            logging.error(f"Presentation file operation failed: {str(e)}")
            return False

    logging.info(f"Starting presentation generation for: {topic}")
    clean_topic = topic.strip()[:100]
    content = generate_presentation_content(clean_topic)
    
    if content.startswith("Error:"):
        logging.error(content)
        return False
        
    filename = clean_topic.lower().replace(" ", "_").replace(":", "").replace("/", "_")
    return save_and_open_presentation(content, filename)

def YoutubeSearch(query: str) -> bool:
    """Search YouTube"""
    try:
        url = f"https://www.youtube.com/results?search_query={query.replace(' ', '+')}"
        webbrowser.open(url)
        logging.info(f"YouTube search executed for: {query}")
        return True
    except Exception as e:
        logging.error(f"YouTube search failed: {str(e)}")
        return False

def PlayYoutube(query: str) -> bool:
    """Play YouTube video"""
    try:
        playonyt(query)
        logging.info(f"YouTube playback started for: {query}")
        return True
    except Exception as e:
        logging.error(f"YouTube playback failed: {str(e)}")
        return False

def OpenApp(app_name: str) -> bool:
    """Open application"""
    try:
        appopen(app_name, match_closest=True, throw_error=True)
        logging.info(f"Opened app: {app_name}")
        return True
    except Exception:
        if not session:
            logging.error(f"Cannot open website for {app_name}: HTTP session not initialized")
            return False
        try:
            url = f"https://www.google.com/search?q={app_name}+official+site"
            response = session.get(url)
            soup = BeautifulSoup(response.text, "html.parser")
            link = soup.find("a", {"jsname": "UWckNb"})
            if link and (href := link.get("href")):
                webbrowser.open(href)
                logging.info(f"Opened website for: {app_name}")
                return True
            logging.error(f"No website found for: {app_name}")
            return False
        except Exception as e:
            logging.error(f"App opening failed: {str(e)}")
            return False

def CloseApp(app_name: str) -> bool:
    """Close application or folder"""
    logging.info(f"Attempting to close: {app_name}")
    try:
        if "youtube" in app_name.lower():
            try:
                close("youtube", match_closest=True, throw_error=True)
                logging.info("YouTube app closed successfully")
                return True
            except Exception:
                logging.info("YouTube app not found, attempting to close browser tab/process")
                try:
                    pyautogui.hotkey("ctrl", "w")  # Close active browser tab
                    logging.info("Closed active browser tab")
                except Exception as e:
                    logging.error(f"Failed to close browser tab: {str(e)}")
                subprocess.run(["taskkill", "/IM", "chrome.exe", "/F"], check=False)
                subprocess.run(["taskkill", "/IM", "firefox.exe", "/F"], check=False)
                subprocess.run(["taskkill", "/IM", "edge.exe", "/F"], check=False)
                logging.info("Terminated browser processes")
                return True
        elif "chrome" in app_name.lower():
            subprocess.run(["taskkill", "/IM", "chrome.exe", "/F"], check=True)
            logging.info("Chrome closed successfully")
            return True
        elif os.path.isdir(app_name):
            try:
                subprocess.run(["taskkill", "/IM", "explorer.exe", "/F"], check=False)
                subprocess.run(["start", "explorer.exe"], shell=True, check=False)
                logging.info(f"Closed folder window: {app_name}")
                return True
            except Exception as e:
                logging.error(f"Failed to close folder: {str(e)}")
                return False
        else:
            close(app_name, match_closest=True, throw_error=True)
            logging.info(f"App {app_name} closed successfully")
            return True
    except Exception as e:
        logging.error(f"App closing failed for {app_name}: {str(e)}")
        return False

def System(command: str) -> bool:
    """System controls"""
    commands = {
        "mute": "volume mute",
        "unmute": "volume mute",
        "volume up": "volume up",
        "volume down": "volume down",
        "shutdown": "shutdown"
    }
    
    try:
        cmd = command.lower().strip()
        if cmd == "shutdown":
            os.system("shutdown /s /t 1" if os.name == "nt" else "shutdown -h now")
            return True
        if cmd in commands:
            keyboard.press_and_release(commands[cmd])
            return True
        return False
    except Exception as e:
        logging.error(f"System command failed: {str(e)}")
        return False





def SendMail(to_address: str, subject: str, body: str) -> bool:
    """Send email"""
    try:
        msg = MIMEMultipart()
        msg["From"] = EMAIL_ADDRESS
        msg["To"] = to_address
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain"))
        
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()
            server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
            server.send_message(msg)
        logging.info(f"Email sent to: {to_address}")
        return True
    except Exception as e:
        logging.error(f"Email sending failed: {str(e)}")
        return False

async def TranslateAndExecute(commands: list[str]) -> list[bool]:
    """Process commands asynchronously"""
    tasks = []
    
    for command in commands:
        cmd = command.strip().lower()
        if not cmd:
            continue
            
        try:
            if cmd.startswith("open "):
                app = cmd[5:].strip()
                tasks.append(asyncio.to_thread(OpenApp, app))
            elif cmd.startswith("close "):
                app = cmd[6:].strip()
                tasks.append(asyncio.to_thread(CloseApp, app))
            elif cmd.startswith("play "):
                query = cmd[5:].strip()
                tasks.append(asyncio.to_thread(PlayYoutube, query))
            elif cmd.startswith(("write ", "content ")):
                topic = cmd.split(" ", 1)[1].strip()
                if not client:
                    logging.error(f"Skipping content generation for '{topic}': Groq client not initialized")
                    continue
                tasks.append(asyncio.to_thread(Content, topic))
            elif cmd.startswith("create presentation "):
                topic = cmd[19:].strip()
                if not client:
                    logging.error(f"Skipping presentation generation for '{topic}': Groq client not initialized")
                    continue
                tasks.append(asyncio.to_thread(CreateGammaPresentation, topic))
            elif cmd.startswith("google search "):
                query = cmd[14:].strip()
                tasks.append(asyncio.to_thread(GoogleSearch, query))
            elif cmd.startswith("youtube search "):
                query = cmd[15:].strip()
                tasks.append(asyncio.to_thread(YoutubeSearch, query))
            elif cmd.startswith("system "):
                command = cmd[7:].strip()
                tasks.append(asyncio.to_thread(System, command))

            elif cmd.startswith("send mail "):
                parts = cmd[9:].split(" about ", 1)
                to_address = parts[0].replace("to ", "").strip()
                subject_body = parts[1].split(" with ", 1) if len(parts) > 1 else ["", ""]
                tasks.append(asyncio.to_thread(
                    SendMail,
                    to_address,
                    subject_body[0].strip(),
                    subject_body[1].strip() if len(subject_body) > 1 else ""
                ))
            else:
                logging.warning(f"Unknown command: {cmd}")
        except Exception as e:
            logging.error(f"Command processing error: {str(e)}")

    results = await asyncio.gather(*tasks, return_exceptions=True)
    return [not isinstance(r, Exception) and r for r in results]

async def Automation(commands: list[str]) -> bool:
    """Main automation controller"""
    logging.info("Starting automation sequence")
    results = await TranslateAndExecute(commands)

    if not results:
        logging.error("No commands executed")
        return False

    success_rate = sum(results) / len(results)
    logging.info(f"Automation completed with {success_rate:.0%} success rate")
    return all(results)

if __name__ == "__main__":
    # Example usage - Basic Automation
    commands = [
        "open notepad",
        "close chrome",
        "system volume up",
        "google search python tutorial"
    ]

    async def main():
        await Automation(commands)

    asyncio.run(main())